#!/usr/bin/env python3
"""
Convert MailingPigeon webpage to PowerPoint presentation
"""

import os
from pathlib import Path
from bs4 import BeautifulSoup

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call(['pip3', 'install', 'python-pptx'])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

# Base directory
BASE_DIR = Path(__file__).parent

def extract_text(element):
    """Extract clean text from HTML element"""
    if element is None:
        return ""
    # Remove script and style tags
    for script in element(["script", "style", "svg"]):
        script.decompose()
    text = element.get_text(separator=" ", strip=True)
    return " ".join(text.split())

def extract_sections(html_file):
    """Parse HTML and extract sections"""
    with open(html_file, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    
    sections = []
    for section in soup.find_all('section', class_='section'):
        section_id = section.get('id', '')
        section_data = {
            'id': section_id,
            'title': '',
            'text': [],
            'images': [],
            'background_image': None
        }
        
        # Extract title (h1, h2, or first large text)
        title_elem = section.find(['h1', 'h2', 'h3'])
        if title_elem:
            section_data['title'] = extract_text(title_elem)
        
        # Extract background image from style attribute
        style = section.get('style', '')
        if 'background-image' in style:
            import re
            match = re.search(r"url\('([^']+)'\)", style)
            if match:
                section_data['background_image'] = match.group(1)
        
        # Extract all images
        for img in section.find_all('img'):
            src = img.get('src', '')
            if src and not img.get('aria-hidden') == 'true':
                alt = img.get('alt', '')
                section_data['images'].append({
                    'src': src,
                    'alt': alt
                })
        
        # Extract text content
        text_elements = section.find_all(['p', 'span', 'div'], class_=lambda x: x and any(
            cls in x for cls in ['story-line', 'story-text', 'mission-line', 'principle', 
                                'transition-line', 'failure-text', 'comparison-line', 
                                'response-line', 'strategy-line', 'finale-line', 'salute-line',
                                'reason-line', 'hospital-line', 'heroic-line', 'mantra-line',
                                'audience-text', 'card-title', 'card-text']
        ))
        
        for elem in text_elements:
            text = extract_text(elem)
            if text and len(text) > 3:
                section_data['text'].append(text)
        
        # If no text found, try getting all paragraph text
        if not section_data['text']:
            for p in section.find_all('p'):
                text = extract_text(p)
                if text and len(text) > 3:
                    section_data['text'].append(text)
        
        if section_data['title'] or section_data['text'] or section_data['images']:
            sections.append(section_data)
    
    return sections

def create_presentation(sections, output_file):
    """Create PowerPoint presentation from sections"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors (WW1 sepia palette)
    sepia_dark = RGBColor(61, 40, 23)
    sepia_brown = RGBColor(92, 64, 51)
    sepia_medium = RGBColor(139, 111, 71)
    
    for i, section in enumerate(sections):
        # Skip empty sections
        if not section['title'] and not section['text'] and not section['images']:
            continue
        
        # Create slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add background image if available
        if section['background_image']:
            bg_path = BASE_DIR / section['background_image']
            if bg_path.exists():
                try:
                    slide.shapes.add_picture(str(bg_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
                except:
                    pass
        
        # Add title
        if section['title']:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = section['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = sepia_dark
            title_para.alignment = PP_ALIGN.LEFT
        
        # Add text content
        text_top = Inches(1.8) if section['title'] else Inches(1)
        for j, text in enumerate(section['text'][:6]):  # Limit to 6 text items per slide
            left = Inches(0.5)
            top = text_top + (j * Inches(0.8))
            width = Inches(6)
            height = Inches(0.7)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = text
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(18)
            text_para.font.color.rgb = sepia_brown
            text_para.alignment = PP_ALIGN.LEFT
        
        # Add images
        img_left = Inches(7)
        img_top = Inches(1.5)
        for k, img_data in enumerate(section['images'][:2]):  # Limit to 2 images per slide
            img_path = BASE_DIR / img_data['src']
            if img_path.exists():
                try:
                    # Calculate size (max 3 inches width)
                    slide.shapes.add_picture(
                        str(img_path),
                        img_left,
                        img_top + (k * Inches(3)),
                        width=Inches(2.5)
                    )
                except Exception as e:
                    print(f"Could not add image {img_data['src']}: {e}")
        
        # Special handling for capability cards
        if section['id'] == 'capabilities':
            # Create a second slide for the 6 capability cards
            card_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Add title
            title_box = card_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = "MailingPigeon does not just deliver emails."
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = sepia_dark
            
            # Add cards in 3x2 grid
            cards = [
                ("🔍", "It scouts.", "Finds the right prospects."),
                ("📋", "It builds.", "Clean, intelligent lists."),
                ("🛤️", "It prepares.", "Warms domains. Warms inboxes."),
                ("📧", "It carries.", "Every kind of message."),
                ("👁️", "It watches.", "Open rates. Replies. Bounces."),
                ("🔄", "It adapts.", "When conditions change.")
            ]
            
            card_width = Inches(2.8)
            card_height = Inches(2)
            card_spacing_x = Inches(0.2)
            card_spacing_y = Inches(0.2)
            start_x = Inches(0.5)
            start_y = Inches(1.5)
            
            for idx, (icon, title, desc) in enumerate(cards):
                row = idx // 3
                col = idx % 3
                x = start_x + (col * (card_width + card_spacing_x))
                y = start_y + (row * (card_height + card_spacing_y))
                
                # Add card background (rectangle)
                from pptx.enum.shapes import MSO_SHAPE
                card_shape = card_slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, card_width, card_height
                )
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = RGBColor(245, 241, 232)  # sepia-cream
                card_shape.line.color.rgb = sepia_medium
                card_shape.line.width = Pt(3)
                
                # Add icon
                icon_box = card_slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.5))
                icon_frame = icon_box.text_frame
                icon_frame.text = icon
                icon_para = icon_frame.paragraphs[0]
                icon_para.font.size = Pt(36)
                icon_para.alignment = PP_ALIGN.CENTER
                
                # Add title
                title_box = card_slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), Inches(2.4), Inches(0.4))
                title_frame = title_box.text_frame
                title_frame.text = title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(20)
                title_para.font.bold = True
                title_para.font.color.rgb = sepia_dark
                title_para.alignment = PP_ALIGN.CENTER
                
                # Add description
                desc_box = card_slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.3), Inches(2.4), Inches(0.5))
                desc_frame = desc_box.text_frame
                desc_frame.text = desc
                desc_para = desc_frame.paragraphs[0]
                desc_para.font.size = Pt(14)
                desc_para.font.color.rgb = sepia_medium
                desc_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save(output_file)
    print(f"✅ Presentation saved to: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    html_file = BASE_DIR / 'index.html'
    output_file = BASE_DIR / 'MailingPigeon_Presentation.pptx'
    
    print("📄 Parsing HTML...")
    sections = extract_sections(html_file)
    print(f"📑 Found {len(sections)} sections")
    
    print("🎨 Creating PowerPoint presentation...")
    create_presentation(sections, output_file)
    
    print("\n✨ Done! Open the .pptx file in PowerPoint or Google Slides.")
